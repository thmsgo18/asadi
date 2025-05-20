from langchain.schema import Document
from .extractions import extraireTextPlumber, extraireTableauxPdf, extraireXLSX, extraire_text_md_Et_txt, extraireDOCX, get_text_from_any_pdf
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from pptx import Presentation
import csv
"""
Module de chargement de fichiers de différents formats en documents Langchain.

Ce module fournit des fonctions pour charger différents types de fichiers (PDF, DOCX, XLSX, etc.)
et les convertir en objets Document de Langchain pour une utilisation ultérieure
dans des systèmes de RAG ou d'autres applications de traitement de texte.
"""

def docxLoad(lienDocx):
    """
    Charge un fichier Word (DOCX) et le convertit en Document Langchain.
    
    Args:
        lienDocx (str): Chemin vers le fichier DOCX à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées.
    """    
    texte = extraireDOCX(lienDocx)
    document = Document(
        page_content=texte,
        metadata={"source": lienDocx}
    )
    return document


def pdfScanneLoad(lienPdf):
    """
    Extrait le texte d'un PDF scanné et le retourne sous le format Document Langchain.
    
    Utilise l'OCR pour extraire le texte des images contenues dans le PDF scanné.
    
    Args:
        lienPdf (str): Chemin vers le fichier PDF scanné à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées.
    """ 
    texte= get_text_from_any_pdf(lienPdf)
    document = Document(
            page_content=texte,
            metadata={"source": lienPdf}
        )
    return document


def excelLoad(lienExcel):
    """
    Charge un fichier Excel (XLSX) et le convertit en Document Langchain.
    
    Extrait le contenu de toutes les feuilles du fichier Excel et inclut
    les noms des feuilles dans les métadonnées.
    
    Args:
        lienExcel (str): Chemin vers le fichier Excel à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées.
    """ 
    texte, nomsFeuilles= extraireXLSX(lienExcel)
    document = Document(
        page_content = texte,
        metadata = {"source" : lienExcel, "nomsFeuilles": ", ".join(nomsFeuilles)}
    )

    return document

def mdLoad(lienMd):
    """
    Charge un fichier Markdown et le convertit en Document Langchain.
    
    Args:
        lienMd (str): Chemin vers le fichier Markdown à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées.
    """
    texte= extraire_text_md_Et_txt(lienMd)
    document = Document(
            page_content=texte,
            metadata={"source": lienMd}
        )
    return document



def pdfNonScanneLoad(lienPdf):
    """
    Extrait le texte d'un PDF non scanné et le retourne sous le format Document Langchain.
    
    Combine le texte principal et les tableaux extraits du PDF pour créer
    un document complet avec toutes les informations disponibles.
    
    Args:
        lienPdf (str): Chemin vers le fichier PDF non scanné à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte et les tableaux extraits.
    """    
    # extraction du texte avec pdfplumber
    texte = extraireTextPlumber(lienPdf)
    # extraction des tableaux et conversion en texte 
    tableaux = extraireTableauxPdf(lienPdf)
    # fusionner les texte et  les tableaux
    contenu_final = texte + "\n\n" + tableaux if tableaux else texte
    document = Document(
        page_content=contenu_final,
        metadata={"source": lienPdf}
    )

    return document

def rtfLoad(lienRtf):
    """
    Charge un fichier RTF et le convertit en Document Langchain.
    
    Utilise la bibliothèque striprtf pour convertir le contenu RTF en texte brut.
    
    Args:
        lienRtf (str): Chemin vers le fichier RTF à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées.
    """

    with open(lienRtf, "r", encoding="utf-8") as file:
        rtf_content = file.read()
    texte = rtf_to_text(rtf_content)

    document = Document(
        page_content=texte,
        metadata={"source": lienRtf}
    )
    return document

def htmlLoad(lienHtml):
    """
    Charge un fichier HTML et le convertit en Document Langchain.
    
    Utilise BeautifulSoup pour extraire le texte du contenu HTML en ignorant les balises.
    
    Args:
        lienHtml (str): Chemin vers le fichier HTML à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées.
    """
    with open(lienHtml, "r", encoding="utf-8") as file:
        html_content = file.read()
    soup = BeautifulSoup(html_content, "html.parser")
    texte = soup.get_text(separator="\n", strip=True)

    document = Document(
        page_content=texte,
        metadata={"source": lienHtml}
    )
    return document


def txtLoad(lienTxt):
    """
    Charge un fichier texte (TXT) et le convertit en Document Langchain.
    
    Args:
        lienTxt (str): Chemin vers le fichier texte à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées.
    """
    texte = extraire_text_md_Et_txt(lienTxt)
    document = Document(
        page_content=texte,
        metadata={"source": lienTxt}
    )
    return document

def csvLoad(lienCsv):
    """
    Charge un fichier CSV et le convertit en Document Langchain.
    
    Lit le fichier CSV ligne par ligne et convertit chaque ligne en texte
    avec les valeurs séparées par des virgules.
    
    Args:
        lienCsv (str): Chemin vers le fichier CSV à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées.
    """
    contenu = []
    with open(lienCsv, newline='', encoding="utf-8") as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            contenu.append(", ".join(row))
    texte = "\n".join(contenu)

    document = Document(
        page_content=texte,
        metadata={"source": lienCsv}
    )
    return document
def pptxLoad(lienPptx):
    """
    Charge un fichier PowerPoint (PPTX) et le convertit en Document Langchain.
    
    Extrait le texte de chaque diapositive et le formate avec des séparateurs
    pour préserver la structure de la présentation.
    
    Args:
        lienPptx (str): Chemin vers le fichier PowerPoint à charger.
        
    Returns:
        Document: Objet Document Langchain contenant le texte extrait et les métadonnées,
                 ou un document vide en cas d'erreur.
    """
    try:
        presentation = Presentation(lienPptx)
        all_slides = []
        for idx, slide in enumerate(presentation.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                all_slides.append(f"--- Slide {idx + 1} ---\n" + "\n".join(slide_texts))
        
        texte = "\n\n".join(all_slides)

        document = Document(
            page_content=texte,
            metadata={"source": lienPptx}
        )
        return document
    except Exception as e:
        print(f"Erreur lors de l'extraction du fichier PPTX {lienPptx} : {e}")
        return Document(
            page_content="",
            metadata={"source": lienPptx}
        )